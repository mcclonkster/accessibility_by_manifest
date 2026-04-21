"""
Script 1 — PPTX extraction and preservation for the DOCX remediation workflow.

Purpose
- Extract workflow-relevant content from one PowerPoint deck.
- Preserve source-side provenance, structural hints, role distinctions, and warnings.
- Produce a remediation-prep DOCX draft.
- Produce an authoritative JSON manifest plus YAML and Markdown mirrors.

Authoritative process truth
- The JSON manifest is the authoritative machine-readable output.
- The YAML file mirrors the JSON manifest exactly.
- The Markdown extract report is a rendered human-readable view of the same data.
- The DOCX is a remediation draft, not a compliance-preserving final artifact.

Outputs
- Companion DOCX
- Review notes Markdown
- Authoritative JSON manifest
- YAML mirror of the manifest
- Markdown extract report

Non-goals
- No OCR
- Slide image rendering is optional and uses local LibreOffice plus pdftoppm when PREVIEW_IMAGE_DIR is blank.
- No batch processing
- No final accessibility judgment
- No faithful preservation of all layout meaning, visual relationships, or Word list semantics

Dependencies
- python-pptx
- python-docx
- pyyaml

Basic usage
- Edit the CONFIG block below.
- Set INPUT_PATH to an existing .pptx file.
- Optionally set PREVIEW_IMAGE_DIR to a folder of slide images.
- Leave output paths blank to derive them automatically from the input stem.
- Run the script directly. No CLI arguments are required.

Manual-review expectations
- Reading order is heuristic.
- Only charts and pictures become visual entries in v1.
- Grouped content, standalone text boxes, connector-based meaning, and weak visual metadata require manual review.
- The review-notes file is part of the intended workflow, not an optional extra.
"""

from __future__ import annotations

import json
import logging
import re
import subprocess
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Iterator, Optional

try:
    import yaml  # type: ignore
except Exception:
    yaml = None

from docx import Document
from docx.enum.text import WD_BREAK
from docx.oxml.ns import qn
from docx.shared import Inches, RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


# ==========================
# CONFIG — CORE INPUT/OUTPUT
# ==========================
INPUT_PATH = "/Users/computerk/Library/CloudStorage/Dropbox/0_VSCode/extracted_chat_artifacts/test_inputs/Disorders_in_the_DSM5TR_Presentation.pptx"

OUTPUT_DOCX_PATH = "/Users/computerk/Library/CloudStorage/Dropbox/0_VSCode/extracted_chat_artifacts/test_outputs/Disorders_in_the_DSM5TR_Presentation_companion.docx"
OUTPUT_JSON_MANIFEST_PATH = "/Users/computerk/Library/CloudStorage/Dropbox/0_VSCode/extracted_chat_artifacts/test_outputs/Disorders_in_the_DSM5TR_Presentation_manifest.json"
OUTPUT_YAML_MANIFEST_PATH = "/Users/computerk/Library/CloudStorage/Dropbox/0_VSCode/extracted_chat_artifacts/test_outputs/Disorders_in_the_DSM5TR_Presentation_manifest.yaml"
OUTPUT_EXTRACT_REPORT_PATH = "/Users/computerk/Library/CloudStorage/Dropbox/0_VSCode/extracted_chat_artifacts/test_outputs/Disorders_in_the_DSM5TR_Presentation_extract_report.md"
OUTPUT_REVIEW_NOTES_PATH = "/Users/computerk/Library/CloudStorage/Dropbox/0_VSCode/extracted_chat_artifacts/test_outputs/Disorders_in_the_DSM5TR_Presentation_review_notes.md"

OVERWRITE = False
DRY_RUN = False

# ==========================
# CONFIG — PREVIEW IMAGES
# ==========================
PREVIEW_IMAGE_DIR = ""
PREVIEW_IMAGE_GLOB = "*"
INCLUDE_PREVIEW_IMAGES = True
AUTO_RENDER_SLIDE_IMAGES = True
SLIDE_IMAGE_RENDER_DPI = 160
LIBREOFFICE_EXECUTABLE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
PDFTOPPM_EXECUTABLE = "/opt/homebrew/bin/pdftoppm"

# ==========================
# CONFIG — DOCUMENT / OUTPUT
# ==========================
DOC_TITLE = ""
SLIDE_SEPARATOR_MODE = "PAGE_BREAK"  # "PAGE_BREAK" | "BLANK_LINE"

# ==========================
# CONFIG — EXTRACTION BEHAVIOR
# ==========================
TITLE_FALLBACK_MODE = "FIRST_TEXT_LINE"  # "FIRST_TEXT_LINE" | "UNTITLED_ONLY"
TEXT_ORDER_MODE = "TOP_LEFT"  # "TOP_LEFT"
INCLUDE_EMPTY_SLIDE_TEXT_SECTION = True
INCLUDE_EMPTY_VISUAL_DESCRIPTIONS_SECTION = False

# ==========================
# CONFIG — WARNINGS / REVIEW
# ==========================
WARN_ON_MISSING_ALT_TEXT = True
WARN_ON_UNSUPPORTED_SHAPES = True
WRITE_REVIEW_NOTES = True

# ==========================
# CONFIG — LOGGING
# ==========================
LOG_LEVEL = "INFO"  # "DEBUG" | "INFO" | "WARNING" | "ERROR"
LOG_TO_FILE = False
LOG_FILE_PATH = "./script1_extract.log"
SHOW_PROGRESS = True

# Input and output rules
# - v1 expects INPUT_PATH to be a single .pptx file.
# - If any output path is blank, it is derived from the input stem.
# - The JSON manifest is the authoritative machine-readable output.
# - The YAML file mirrors the JSON manifest exactly.
# - The Markdown extract report is a rendered human-readable view of the same data.
# - The DOCX is a remediation draft and not a compliance-preserving final artifact.
# - Preview-image candidates in the manifest are written as relative filenames/paths.
# - When DRY_RUN=True, the script performs validation, extraction, manifest construction,
#   and reporting, but makes zero filesystem changes.


EXIT_SUCCESS = 0
EXIT_CONFIG_ERROR = 2
EXIT_RUNTIME_ERROR = 3

STATUS_SUCCESS = "SUCCESS"
STATUS_SUCCESS_WITH_WARNINGS = "SUCCESS_WITH_WARNINGS"
STATUS_DRY_RUN_COMPLETE = "DRY_RUN_COMPLETE"
STATUS_FAILURE = "FAILURE"

ALLOWED_LOG_LEVELS = {"DEBUG", "INFO", "WARNING", "ERROR"}
ALLOWED_TITLE_FALLBACK_MODES = {"FIRST_TEXT_LINE", "UNTITLED_ONLY"}
ALLOWED_TEXT_ORDER_MODES = {"TOP_LEFT"}
ALLOWED_SLIDE_SEPARATOR_MODES = {"PAGE_BREAK", "BLANK_LINE"}

SOURCE_TEXT_CONTAINER_TYPES = {"placeholder_text", "standalone_text_box", "table_cell"}
PLACEHOLDER_TYPES = {"title", "body", "subtitle", "other"}
INTERPRETED_TEXT_ROLES = {
    "title_candidate",
    "body_content",
    "caption",
    "callout",
    "label",
    "annotation",
    "unknown",
}
ROLE_CONFIDENCE_VALUES = {"high", "medium", "low"}
LIST_KIND_VALUES = {"plain", "bullet", "number"}
DESCRIPTION_SOURCE_VALUES = {"descr", "title", "missing"}
PROJECTION_STATUS_VALUES = {"planned", "written", "omitted", "manual_only"}

SUPPORTED_PREVIEW_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}
PREVIEW_NUMBER_PATTERN = re.compile(r"(?i)(?:^|[^0-9])(?:slide)?[_\-\s]*0*([1-9][0-9]*)(?:[^0-9]|$)")
XML_ILLEGAL_CHAR_PATTERN = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\ud800-\udfff]")

SECTION_HEADING_SLIDE_TEXT = "Slide text"
SECTION_HEADING_VISUALS = "Visual descriptions"

PLACEHOLDER_MISSING_TEXT = "No extractable visible slide text found."
PLACEHOLDER_MISSING_DESCRIPTION = "Description metadata missing. Review needed."

EMU_PER_INCH = 914400
READING_ORDER_VERTICAL_THRESHOLD = int(0.35 * EMU_PER_INCH)
READING_ORDER_HORIZONTAL_THRESHOLD = int(2.5 * EMU_PER_INCH)
NONTRIVIAL_SHAPE_MIN_SIZE = int(0.3 * EMU_PER_INCH)
VISUAL_HEAVY_TEXT_BLOCK_THRESHOLD = 1
LAYOUT_CONTEXT_DISTANCE_THRESHOLD = int(1.5 * EMU_PER_INCH)

TIMESTAMP_UTC = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

LOGGER = logging.getLogger("script1_extract")
WCAG_AAA_TEXT_COLOR = RGBColor(0, 0, 0)
WCAG_AAA_TEXT_COLOR_HEX = "000000"


class ConfigError(Exception):
    """Raised when CONFIG values are invalid or incompatible."""


class SimpleShapeProxy:
    """Minimal proxy used to carry table-cell geometry and naming into shared logic."""

    def __init__(self, name: str, left: int, top: int, width: int, height: int, shape_id: Optional[int] = None) -> None:
        self.name = name
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.shape_id = shape_id


class OutputPaths:
    def __init__(
        self,
        docx: Path,
        manifest_json: Path,
        manifest_yaml: Path,
        extract_report_md: Path,
        review_notes_md: Path,
    ) -> None:
        self.docx = docx
        self.manifest_json = manifest_json
        self.manifest_yaml = manifest_yaml
        self.extract_report_md = extract_report_md
        self.review_notes_md = review_notes_md


class ShapeContext:
    def __init__(self, shape: Any, inside_group: bool, is_group_container: bool) -> None:
        self.shape = shape
        self.inside_group = inside_group
        self.is_group_container = is_group_container


def setup_logging() -> logging.Logger:
    logger = LOGGER
    logger.handlers.clear()
    logger.setLevel(getattr(logging, LOG_LEVEL, logging.INFO))
    logger.propagate = False

    formatter = logging.Formatter("%(asctime)s | %(levelname)s | %(message)s")

    stream_handler = logging.StreamHandler(sys.stderr)
    stream_handler.setFormatter(formatter)
    logger.addHandler(stream_handler)

    if LOG_TO_FILE and not DRY_RUN:
        log_path = Path(LOG_FILE_PATH)
        log_path.parent.mkdir(parents=True, exist_ok=True)
        file_handler = logging.FileHandler(log_path, encoding="utf-8")
        file_handler.setFormatter(formatter)
        logger.addHandler(file_handler)

    return logger


def fail(message: str, exit_code: int) -> int:
    LOGGER.error(message)
    return exit_code


def derive_output_paths(input_path: Path) -> OutputPaths:
    stem = input_path.stem
    parent = input_path.parent

    docx = Path(OUTPUT_DOCX_PATH).expanduser() if OUTPUT_DOCX_PATH else parent / f"{stem}_companion.docx"
    manifest_json = (
        Path(OUTPUT_JSON_MANIFEST_PATH).expanduser()
        if OUTPUT_JSON_MANIFEST_PATH
        else parent / f"{stem}_extract_manifest.json"
    )
    manifest_yaml = (
        Path(OUTPUT_YAML_MANIFEST_PATH).expanduser()
        if OUTPUT_YAML_MANIFEST_PATH
        else parent / f"{stem}_extract_manifest.yaml"
    )
    extract_report_md = (
        Path(OUTPUT_EXTRACT_REPORT_PATH).expanduser()
        if OUTPUT_EXTRACT_REPORT_PATH
        else parent / f"{stem}_extract_report.md"
    )
    review_notes_md = (
        Path(OUTPUT_REVIEW_NOTES_PATH).expanduser()
        if OUTPUT_REVIEW_NOTES_PATH
        else parent / f"{stem}_review_notes.md"
    )

    return OutputPaths(
        docx=docx,
        manifest_json=manifest_json,
        manifest_yaml=manifest_yaml,
        extract_report_md=extract_report_md,
        review_notes_md=review_notes_md,
    )


def validate_config() -> tuple[Path, OutputPaths, Optional[Path]]:
    if yaml is None:
        raise ConfigError("PyYAML is required. Install it with: pip install pyyaml")

    input_path = Path(INPUT_PATH).expanduser()
    if not INPUT_PATH:
        raise ConfigError("INPUT_PATH is empty. Set INPUT_PATH to an existing .pptx file.")
    if not input_path.exists():
        raise ConfigError("INPUT_PATH does not exist. Set INPUT_PATH to an existing .pptx file.")
    if not input_path.is_file() or input_path.suffix.lower() != ".pptx":
        raise ConfigError("INPUT_PATH must point to a .pptx file, not a directory or unsupported file type.")

    if TITLE_FALLBACK_MODE not in ALLOWED_TITLE_FALLBACK_MODES:
        allowed = ", ".join(sorted(ALLOWED_TITLE_FALLBACK_MODES))
        raise ConfigError(f"TITLE_FALLBACK_MODE must be one of: {allowed}.")
    if TEXT_ORDER_MODE not in ALLOWED_TEXT_ORDER_MODES:
        allowed = ", ".join(sorted(ALLOWED_TEXT_ORDER_MODES))
        raise ConfigError(f"TEXT_ORDER_MODE must be one of: {allowed}.")
    if SLIDE_SEPARATOR_MODE not in ALLOWED_SLIDE_SEPARATOR_MODES:
        allowed = ", ".join(sorted(ALLOWED_SLIDE_SEPARATOR_MODES))
        raise ConfigError(f"SLIDE_SEPARATOR_MODE must be one of: {allowed}.")
    if LOG_LEVEL not in ALLOWED_LOG_LEVELS:
        allowed = ", ".join(sorted(ALLOWED_LOG_LEVELS))
        raise ConfigError(f"LOG_LEVEL must be one of: {allowed}.")

    output_paths = derive_output_paths(input_path)
    for path in [
        output_paths.docx,
        output_paths.manifest_json,
        output_paths.manifest_yaml,
        output_paths.extract_report_md,
        output_paths.review_notes_md,
    ]:
        if path.suffix == "":
            raise ConfigError(f"Output path must be a file path, not a directory: {path}")
        if path.exists() and path.is_dir():
            raise ConfigError(f"Output path points to an existing directory: {path}")
        if path.parent.exists() and not path.parent.is_dir():
            raise ConfigError(f"Output path parent exists but is not a directory: {path.parent}")
        if path.exists() and not OVERWRITE and not DRY_RUN:
            raise ConfigError(f"Output file already exists. Set OVERWRITE = True or change the path: {path}")

    preview_dir: Optional[Path] = None
    if PREVIEW_IMAGE_DIR:
        preview_dir = Path(PREVIEW_IMAGE_DIR).expanduser()
        if not preview_dir.exists() or not preview_dir.is_dir():
            raise ConfigError(
                "PREVIEW_IMAGE_DIR is set but is not an existing directory. "
                "Set it to a valid folder or leave it empty."
            )

    return input_path, output_paths, preview_dir


def ensure_output_directories(output_paths: OutputPaths) -> None:
    for path in [
        output_paths.docx,
        output_paths.manifest_json,
        output_paths.manifest_yaml,
        output_paths.extract_report_md,
        output_paths.review_notes_md,
    ]:
        if path.parent.exists():
            continue
        if DRY_RUN:
            LOGGER.info("DRY_RUN: would create directory %s", path.parent)
        else:
            path.parent.mkdir(parents=True, exist_ok=True)


def load_presentation(input_path: Path) -> Any:
    try:
        presentation = Presentation(str(input_path))
    except Exception as exc:
        raise RuntimeError(f"Failed to open PowerPoint file '{input_path}': {exc}") from exc
    LOGGER.info("Loaded presentation with %s slide(s).", len(presentation.slides))
    return presentation


def render_slide_images_if_needed(input_path: Path, output_paths: OutputPaths, preview_dir: Optional[Path]) -> Optional[Path]:
    if not INCLUDE_PREVIEW_IMAGES:
        return preview_dir
    if preview_dir is not None:
        return preview_dir
    if not AUTO_RENDER_SLIDE_IMAGES:
        return None

    rendered_dir = output_paths.docx.parent / f"{input_path.stem}_slide_images"
    if DRY_RUN:
        LOGGER.info("DRY_RUN: would render slide images to %s", rendered_dir)
        return rendered_dir

    render_slide_images(input_path, rendered_dir)

    global PREVIEW_IMAGE_DIR
    PREVIEW_IMAGE_DIR = str(rendered_dir)
    return rendered_dir


def render_slide_images(input_path: Path, rendered_dir: Path) -> None:
    soffice = Path(LIBREOFFICE_EXECUTABLE)
    pdftoppm = Path(PDFTOPPM_EXECUTABLE)
    if not soffice.exists():
        raise RuntimeError(f"LibreOffice executable not found: {soffice}")
    if not pdftoppm.exists():
        raise RuntimeError(f"pdftoppm executable not found: {pdftoppm}")

    rendered_dir.mkdir(parents=True, exist_ok=True)

    for old_image in rendered_dir.glob("slide_*.png"):
        old_image.unlink()

    with tempfile.TemporaryDirectory(prefix="pptx_slide_render_") as temp_root:
        temp_root_path = Path(temp_root)
        pdf_dir = temp_root_path / "pdf"
        profile_dir = temp_root_path / "lo_profile"
        pdf_dir.mkdir(parents=True, exist_ok=True)
        profile_dir.mkdir(parents=True, exist_ok=True)

        convert_cmd = [
            str(soffice),
            f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_dir),
            str(input_path),
        ]
        LOGGER.info("Rendering slide images: converting PPTX to PDF.")
        subprocess.run(convert_cmd, check=True, capture_output=True, text=True)

        pdf_path = pdf_dir / f"{input_path.stem}.pdf"
        if not pdf_path.exists():
            pdf_candidates = sorted(pdf_dir.glob("*.pdf"))
            if not pdf_candidates:
                raise RuntimeError("LibreOffice did not produce a PDF for slide image rendering.")
            pdf_path = pdf_candidates[0]

        image_prefix = temp_root_path / "slide"
        raster_cmd = [
            str(pdftoppm),
            "-r",
            str(SLIDE_IMAGE_RENDER_DPI),
            "-png",
            str(pdf_path),
            str(image_prefix),
        ]
        LOGGER.info("Rendering slide images: rasterizing PDF pages.")
        subprocess.run(raster_cmd, check=True, capture_output=True, text=True)

        rendered_pages = sorted(
            temp_root_path.glob("slide-*.png"),
            key=lambda path: extract_slide_number_from_preview_name(path) or 0,
        )
        if not rendered_pages:
            raise RuntimeError("pdftoppm did not produce slide images.")

        for page in rendered_pages:
            slide_number = extract_slide_number_from_preview_name(page)
            if slide_number is None:
                continue
            page.replace(rendered_dir / f"slide_{slide_number:03d}.png")

    LOGGER.info("Rendered %s slide image(s) to %s.", len(list(rendered_dir.glob("slide_*.png"))), rendered_dir)


def iter_shape_contexts(shapes: Iterable[Any], inside_group: bool = False) -> Iterator[ShapeContext]:
    for shape in shapes:
        is_group_container = getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
        yield ShapeContext(shape=shape, inside_group=inside_group, is_group_container=is_group_container)
        if is_group_container:
            try:
                yield from iter_shape_contexts(shape.shapes, inside_group=True)
            except Exception:
                continue


def discover_preview_images(preview_dir: Optional[Path]) -> tuple[dict[int, Path], dict[int, list[dict[str, Any]]]]:
    if not INCLUDE_PREVIEW_IMAGES or preview_dir is None:
        return {}, {}

    candidates = sorted(preview_dir.glob(PREVIEW_IMAGE_GLOB))
    preview_matches: dict[int, list[Path]] = {}
    for path in candidates:
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_PREVIEW_SUFFIXES:
            continue
        slide_number = extract_slide_number_from_preview_name(path)
        if slide_number is None:
            continue
        preview_matches.setdefault(slide_number, []).append(path)

    preview_lookup: dict[int, Path] = {}
    preview_warnings_by_slide: dict[int, list[dict[str, Any]]] = {}
    for slide_number, matches in preview_matches.items():
        preview_lookup[slide_number] = matches[0]
        if len(matches) > 1:
            preview_warnings_by_slide[slide_number] = [
                make_warning(
                    "PREVIEW_MATCH_AMBIGUOUS",
                    "Multiple preview files matched this slide number; first sorted match used: "
                    + ", ".join(relative_preview_candidate(m) for m in matches),
                    "slide",
                    True,
                )
            ]

    LOGGER.info("Matched %s preview image candidate(s).", len(preview_lookup))
    return preview_lookup, preview_warnings_by_slide


def extract_slide_number_from_preview_name(path: Path) -> Optional[int]:
    match = PREVIEW_NUMBER_PATTERN.search(path.stem)
    if not match:
        return None
    try:
        return int(match.group(1))
    except ValueError:
        return None


def relative_preview_candidate(path: Path) -> str:
    if PREVIEW_IMAGE_DIR:
        base = Path(PREVIEW_IMAGE_DIR).expanduser()
        try:
            return str(path.relative_to(base))
        except Exception:
            return path.name
    return path.name


def safe_shape_name(shape: Any) -> str:
    name = getattr(shape, "name", "") or getattr(shape, "shape_id", "")
    return str(name) if name else "Unnamed shape"


def safe_shape_type_name(shape: Any) -> str:
    try:
        shape_type = getattr(shape, "shape_type", None)
        if shape_type is None:
            return "UNKNOWN"
        return getattr(shape_type, "name", str(shape_type))
    except Exception:
        return "UNKNOWN"


def safe_shape_kind(shape: Any) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "p:pic"
    if getattr(shape, "has_chart", False):
        return "p:graphicFrame"
    if getattr(shape, "is_placeholder", False) or getattr(shape, "has_text_frame", False):
        return "p:sp"
    return safe_shape_type_name(shape)


def safe_shape_id(shape: Any) -> Optional[int]:
    try:
        return int(getattr(shape, "shape_id", 0))
    except Exception:
        return None


def geometry_from_shape(shape: Any) -> Optional[dict[str, int]]:
    try:
        return {
            "x": int(getattr(shape, "left", 0)),
            "y": int(getattr(shape, "top", 0)),
            "cx": int(getattr(shape, "width", 0)),
            "cy": int(getattr(shape, "height", 0)),
        }
    except Exception:
        return None


def make_warning(code: str, message: str, scope: str, manual_review_required: bool) -> dict[str, Any]:
    return {
        "warning_code": code,
        "warning_message": message,
        "warning_scope": scope,
        "manual_review_required": manual_review_required,
    }


def dedupe_warning_entries(entries: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[tuple[str, str, str]] = set()
    output: list[dict[str, Any]] = []
    for entry in entries:
        key = (entry["warning_code"], entry["warning_message"], entry["warning_scope"])
        if key in seen:
            continue
        seen.add(key)
        output.append(entry)
    return output


def normalize_paragraph_text(text: str) -> str:
    if not text:
        return ""
    text = XML_ILLEGAL_CHAR_PATTERN.sub("", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
    lines = [line for line in lines if line]
    return " ".join(lines).strip()


def get_placeholder_semantics(shape: Any) -> tuple[bool, Optional[str]]:
    if not getattr(shape, "is_placeholder", False):
        return False, None

    try:
        placeholder_type = shape.placeholder_format.type
    except Exception:
        return False, None

    if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
        return True, "title"
    if placeholder_type in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}:
        return True, "body"
    if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
        return True, "subtitle"
    return True, "other"


def classify_source_text_container(has_placeholder_semantics: bool, source_kind: str) -> str:
    if source_kind == "table_cell":
        return "table_cell"
    if has_placeholder_semantics:
        return "placeholder_text"
    return "standalone_text_box"


def infer_text_role(
    source_text_container_type: str,
    placeholder_type: Optional[str],
    near_visual: bool,
    relationship_warning_present: bool,
) -> tuple[str, str, bool, list[dict[str, Any]]]:
    warnings: list[dict[str, Any]] = []

    if source_text_container_type == "placeholder_text":
        if placeholder_type == "title":
            return "title_candidate", "high", False, warnings
        if placeholder_type == "body":
            return "body_content", "high", True, warnings
        if placeholder_type == "subtitle":
            return "annotation", "medium", False, warnings
        return "unknown", "medium", False, warnings

    if source_text_container_type == "table_cell":
        return "body_content", "medium", False, warnings

    warnings.append(
        make_warning(
            "STANDALONE_TEXT_BOX_ROLE_UNCERTAIN",
            "Standalone text box role is uncertain.",
            "text_block",
            True,
        )
    )
    warnings.append(
        make_warning(
            "STANDALONE_TEXT_BOX_NOT_MERGED",
            "Standalone text box was preserved separately and not auto-merged.",
            "text_block",
            True,
        )
    )

    if near_visual:
        warnings.append(
            make_warning(
                "STANDALONE_TEXT_BOX_MAY_BE_CAPTION_OR_LABEL",
                "Standalone text box may function as a caption, label, or callout.",
                "text_block",
                True,
            )
        )

    if relationship_warning_present:
        warnings.append(
            make_warning(
                "STANDALONE_TEXT_BOX_MAY_DEPEND_ON_LAYOUT_CONTEXT",
                "Standalone text box may depend on layout context or nearby relationships for full meaning.",
                "text_block",
                True,
            )
        )

    return "unknown", "low", False, warnings


def get_paragraph_level(paragraph: Any) -> int:
    try:
        return int(getattr(paragraph, "level", 0))
    except Exception:
        return 0


def get_list_format_hints(paragraph: Any) -> tuple[str, str]:
    try:
        p = paragraph._p
    except Exception:
        return "plain", ""

    for child in p.iter():
        tag = str(getattr(child, "tag", ""))
        local = tag.rsplit("}", 1)[-1]
        if local == "buAutoNum":
            return "number", "[numbered]"
        if local == "buChar":
            return "bullet", child.get("char", "•") or "•"
        if local == "buBlip":
            return "bullet", "•"
        if local == "buNone":
            return "plain", ""

    return "plain", ""


def list_markup_kind_from_paragraph(paragraph: Any) -> Optional[str]:
    try:
        p = paragraph._p
    except Exception:
        return None

    for child in p.iter():
        tag = str(getattr(child, "tag", ""))
        local = tag.rsplit("}", 1)[-1]
        if local in {"buAutoNum", "buChar", "buBlip", "buNone"}:
            return f"a:{local}"
    return None


def list_markup_value_from_paragraph(paragraph: Any) -> Optional[str]:
    try:
        p = paragraph._p
    except Exception:
        return None

    for child in p.iter():
        tag = str(getattr(child, "tag", ""))
        local = tag.rsplit("}", 1)[-1]
        if local == "buAutoNum":
            return child.get("type", "")
        if local == "buChar":
            return child.get("char", "")
    return None


def paragraph_projection(list_kind: str, ordered_sequence_reconstructable: bool) -> dict[str, Any]:
    if list_kind == "bullet":
        return {
            "target_package_part": "word/document.xml",
            "target_object_kind": "list_paragraph",
            "target_path_hint": "/w:document/w:body",
            "paragraph_style": "ListParagraph",
            "numbering_plan": {"mode": "bullet_list", "num_id": None, "ilvl": None},
            "projection_status": "planned",
        }
    if list_kind == "number":
        if ordered_sequence_reconstructable:
            return {
                "target_package_part": "word/document.xml",
                "target_object_kind": "list_paragraph",
                "target_path_hint": "/w:document/w:body",
                "paragraph_style": "ListParagraph",
                "numbering_plan": {"mode": "ordered_list", "num_id": None, "ilvl": None},
                "projection_status": "planned",
            }
        return {
            "target_package_part": "word/document.xml",
            "target_object_kind": "paragraph",
            "target_path_hint": "/w:document/w:body",
            "paragraph_style": "Normal",
            "numbering_plan": {"mode": "plain_fallback", "num_id": None, "ilvl": None},
            "projection_status": "planned",
        }

    return {
        "target_package_part": "word/document.xml",
        "target_object_kind": "paragraph",
        "target_path_hint": "/w:document/w:body",
        "paragraph_style": "Normal",
        "numbering_plan": {"mode": "plain", "num_id": None, "ilvl": None},
        "projection_status": "planned",
    }


def format_paragraph_for_docx(paragraph_entry: dict[str, Any]) -> str:
    normalized = paragraph_entry["normalized_workflow"]
    text = normalized["text"]
    level = normalized["level"]
    list_kind = normalized["list_kind"]
    marker_hint = normalized["list_marker_hint"]

    indent = "    " * max(level, 0)
    if list_kind == "bullet":
        prefix = f"{marker_hint or '•'} "
    elif list_kind == "number":
        prefix = "[numbered] "
    else:
        prefix = ""
    return f"{indent}{prefix}{text}"


def is_connector_like(shape: Any) -> bool:
    type_name = safe_shape_type_name(shape).upper()
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.LINE:
        return True
    return any(token in type_name for token in ("CONNECTOR", "LINE", "FREEFORM"))


def shape_has_nontrivial_size(shape: Any) -> bool:
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return width >= NONTRIVIAL_SHAPE_MIN_SIZE and height >= NONTRIVIAL_SHAPE_MIN_SIZE


def detect_relationship_meaning_warning(shape_contexts: list[ShapeContext]) -> Optional[dict[str, Any]]:
    connector_like_found = False
    nonpicture_nonchart_shape_count = 0

    for ctx in shape_contexts:
        if ctx.is_group_container:
            continue
        shape = ctx.shape
        if not shape_has_nontrivial_size(shape):
            continue
        if is_connector_like(shape):
            connector_like_found = True
            continue
        if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE and not getattr(shape, "has_chart", False):
            nonpicture_nonchart_shape_count += 1

    if connector_like_found and nonpicture_nonchart_shape_count >= 2:
        return make_warning(
            "CONNECTOR_RELATIONSHIP_WARNING",
            "This slide may use shape relationships or connectors to convey meaning that the output does not preserve.",
            "slide",
            True,
        )
    return None


def should_ignore_shape_for_skip_warning(shape: Any) -> bool:
    if is_connector_like(shape):
        return True
    type_name = safe_shape_type_name(shape).upper()
    if "LINE" in type_name:
        return True
    return False


def shape_is_directly_supported(shape: Any) -> bool:
    if getattr(shape, "has_text_frame", False):
        return True
    if getattr(shape, "has_table", False):
        return True
    if getattr(shape, "has_chart", False):
        return True
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return True
    return False


def detect_skipped_shape_warnings(shape_contexts: list[ShapeContext]) -> list[dict[str, Any]]:
    warnings: list[dict[str, Any]] = []
    if not WARN_ON_UNSUPPORTED_SHAPES:
        return warnings

    for ctx in shape_contexts:
        if ctx.is_group_container:
            continue
        shape = ctx.shape
        if should_ignore_shape_for_skip_warning(shape):
            continue
        if shape_is_directly_supported(shape):
            continue
        if not shape_has_nontrivial_size(shape):
            continue

        warnings.append(
            make_warning(
                "SKIPPED_SHAPE_CONTENT_DETECTED",
                f"Shape '{safe_shape_name(shape)}' of type '{safe_shape_type_name(shape)}' may carry meaning that this extractor does not represent.",
                "slide",
                True,
            )
        )
    return dedupe_warning_entries(warnings)


def is_near_visual(shape: Any, shape_contexts: list[ShapeContext], visual_shapes: set[int]) -> bool:
    shape_left = int(getattr(shape, "left", 0))
    shape_top = int(getattr(shape, "top", 0))

    for ctx in shape_contexts:
        other = ctx.shape
        if id(other) not in visual_shapes:
            continue
        other_left = int(getattr(other, "left", 0))
        other_top = int(getattr(other, "top", 0))
        if abs(shape_left - other_left) <= LAYOUT_CONTEXT_DISTANCE_THRESHOLD and abs(shape_top - other_top) <= LAYOUT_CONTEXT_DISTANCE_THRESHOLD:
            return True
    return False


def extract_paragraph_entries(
    text_frame: Any,
    drawingml_base_path: str,
) -> list[dict[str, Any]]:
    paragraph_entries: list[dict[str, Any]] = []
    for index, paragraph in enumerate(getattr(text_frame, "paragraphs", []), start=1):
        text = normalize_paragraph_text(getattr(paragraph, "text", ""))
        if not text:
            continue

        list_kind, list_marker_hint = get_list_format_hints(paragraph)
        level = get_paragraph_level(paragraph)
        ordered_sequence_reconstructable = False

        paragraph_entries.append(
            {
                "observed_source": {
                    "drawingml_paragraph_path": f"{drawingml_base_path}/a:p[{index}]",
                    "paragraph_level_raw": level,
                    "list_markup_kind": list_markup_kind_from_paragraph(paragraph),
                    "list_markup_value": list_markup_value_from_paragraph(paragraph),
                },
                "normalized_workflow": {
                    "text": text,
                    "level": level,
                    "list_kind": list_kind,
                    "list_marker_hint": list_marker_hint,
                    "ordered_sequence_reconstructable": ordered_sequence_reconstructable,
                },
                "projected_target": paragraph_projection(list_kind, ordered_sequence_reconstructable),
                "warning_entries": (
                    [
                        make_warning(
                            "TRUE_ORDERED_SEQUENCE_NOT_PRESERVED",
                            "Numbering intent was observed, but true ordered sequence is not reconstructable.",
                            "paragraph",
                            True,
                        )
                    ]
                    if list_kind == "number" and not ordered_sequence_reconstructable
                    else []
                ),
            }
        )
    return paragraph_entries


def build_text_block_entry(
    package_part: str,
    shape_tree_path: str,
    shape: Any,
    source_kind: str,
    source_name: str,
    paragraph_entries: list[dict[str, Any]],
    from_group: bool,
    near_visual: bool,
    relationship_warning_present: bool,
) -> dict[str, Any]:
    has_placeholder_semantics, placeholder_type = get_placeholder_semantics(shape)
    source_text_container_type = classify_source_text_container(has_placeholder_semantics, source_kind)
    interpreted_text_role, role_confidence, auto_merge_allowed, role_warnings = infer_text_role(
        source_text_container_type,
        placeholder_type,
        near_visual,
        relationship_warning_present,
    )

    return {
        "observed_source": {
            "package_part": package_part,
            "shape_tree_path": shape_tree_path,
            "shape_kind": safe_shape_kind(shape),
            "shape_id": safe_shape_id(shape),
            "shape_name": safe_shape_name(shape),
            "placeholder_type": placeholder_type,
            "has_placeholder_semantics": has_placeholder_semantics,
            "source_text_container_type": source_text_container_type,
            "geometry": geometry_from_shape(shape),
            "text_body_path": "p:txBody" if getattr(shape, "has_text_frame", False) else None,
            "from_group": from_group,
        },
        "normalized_workflow": {
            "source_kind": source_kind,
            "source_name": source_name,
            "ordered_position": 0,
            "reading_order_sort_key": {
                "top": int(getattr(shape, "top", 0)),
                "left": int(getattr(shape, "left", 0)),
            },
            "interpreted_text_role": interpreted_text_role,
            "role_confidence": role_confidence,
            "auto_merge_allowed": auto_merge_allowed,
        },
        "projected_target": {
            "target_package_part": "word/document.xml",
            "target_object_kind": "paragraph_sequence",
            "target_path_hint": "/w:document/w:body",
            "projection_status": "planned",
        },
        "warning_entries": role_warnings,
        "paragraph_entries": paragraph_entries,
    }


def extract_table_text_blocks(
    package_part: str,
    shape_tree_path: str,
    shape: Any,
    from_group: bool,
    relationship_warning_present: bool,
) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    table = shape.table
    row_count = len(table.rows)
    col_count = len(table.columns)

    for row_idx in range(row_count):
        for col_idx in range(col_count):
            cell = table.cell(row_idx, col_idx)
            paragraph_entries = extract_paragraph_entries(
                cell.text_frame,
                f"{shape_tree_path}/a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]/a:txBody",
            )
            if not paragraph_entries:
                continue

            cell_shape_proxy = SimpleShapeProxy(
                name=f"{safe_shape_name(shape)}[r{row_idx + 1}c{col_idx + 1}]",
                left=int(getattr(shape, "left", 0)) + col_idx,
                top=int(getattr(shape, "top", 0)) + row_idx,
                width=0,
                height=0,
            )
            entry = build_text_block_entry(
                package_part,
                shape_tree_path,
                cell_shape_proxy,
                "table_cell",
                f"{safe_shape_name(shape)}[r{row_idx + 1}c{col_idx + 1}]",
                paragraph_entries,
                from_group,
                False,
                relationship_warning_present,
            )
            blocks.append(entry)

    return blocks


def sort_text_block_entries(text_block_entries: list[dict[str, Any]]) -> list[dict[str, Any]]:
    if TEXT_ORDER_MODE != "TOP_LEFT":
        return text_block_entries

    return sorted(
        text_block_entries,
        key=lambda entry: (
            entry["normalized_workflow"]["reading_order_sort_key"]["top"],
            entry["normalized_workflow"]["reading_order_sort_key"]["left"],
            entry["normalized_workflow"]["source_name"],
        ),
    )


def is_reading_order_weak(text_block_entries: list[dict[str, Any]]) -> bool:
    if len(text_block_entries) < 3:
        return False

    for index, current in enumerate(text_block_entries):
        for later in text_block_entries[index + 1 :]:
            current_top = current["normalized_workflow"]["reading_order_sort_key"]["top"]
            current_left = current["normalized_workflow"]["reading_order_sort_key"]["left"]
            later_top = later["normalized_workflow"]["reading_order_sort_key"]["top"]
            later_left = later["normalized_workflow"]["reading_order_sort_key"]["left"]

            vertical_close = abs(current_top - later_top) <= READING_ORDER_VERTICAL_THRESHOLD
            far_apart_horizontally = abs(current_left - later_left) >= READING_ORDER_HORIZONTAL_THRESHOLD
            if vertical_close and far_apart_horizontally:
                return True
    return False


def flatten_text_block_entries_for_docx(text_block_entries: list[dict[str, Any]]) -> list[str]:
    lines: list[str] = []
    for block in text_block_entries:
        for paragraph_entry in block["paragraph_entries"]:
            lines.append(format_paragraph_for_docx(paragraph_entry))
    return lines


def choose_title(slide: Any, text_block_entries: list[dict[str, Any]], slide_number: int) -> tuple[str, str, list[dict[str, Any]]]:
    warnings: list[dict[str, Any]] = []

    for shape in slide.shapes:
        has_placeholder, placeholder_type = get_placeholder_semantics(shape)
        if has_placeholder and placeholder_type == "title":
            title_text = normalize_paragraph_text(getattr(shape, "text", ""))
            if title_text:
                return title_text, "TITLE_PLACEHOLDER", warnings

    warnings.append(
        make_warning(
            "NO_TITLE_PLACEHOLDER_FOUND",
            "No title placeholder found for this slide.",
            "slide",
            True,
        )
    )

    if TITLE_FALLBACK_MODE == "FIRST_TEXT_LINE":
        for block in text_block_entries:
            paragraphs = block["paragraph_entries"]
            if paragraphs:
                title_text = paragraphs[0]["normalized_workflow"]["text"]
                if title_text:
                    warnings.append(
                        make_warning(
                            "FALLBACK_TITLE_USED",
                            "Fallback title used from extracted text.",
                            "slide",
                            True,
                        )
                    )
                    if len(title_text) > 120:
                        warnings.append(
                            make_warning(
                                "FALLBACK_TITLE_WEAK",
                                "Fallback title came from a long text block and may be weak.",
                                "slide",
                                True,
                            )
                        )
                    return title_text[:200], "FIRST_TEXT_LINE", warnings

    warnings.append(
        make_warning(
            "SAFE_FALLBACK_TITLE_USED",
            "Used safe fallback title because no usable title was found.",
            "slide",
            True,
        )
    )
    return f"Untitled slide {slide_number}", "UNTITLED_FALLBACK", warnings


def extract_visual_entries(
    package_part: str,
    relationships_part: str,
    shape_contexts: list[ShapeContext],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    visual_entries: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    visual_index = 0

    for index, ctx in enumerate(shape_contexts, start=1):
        if ctx.is_group_container:
            continue

        shape = ctx.shape
        if not (getattr(shape, "has_chart", False) or getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE):
            continue

        visual_index += 1
        description, description_source = get_shape_alt_text(shape)
        visual_type, uncertain_type = infer_visual_type(shape, description)
        visual_warnings: list[dict[str, Any]] = []

        if uncertain_type:
            visual_warnings.append(
                make_warning(
                    "VISUAL_TYPE_UNCERTAIN",
                    f"Visual type for '{safe_shape_name(shape)}' is uncertain.",
                    "visual",
                    True,
                )
            )

        if description_source == "title":
            visual_warnings.append(
                make_warning(
                    "DESCRIPTION_SOURCE_TITLE_ONLY",
                    f"Only title metadata was found for visual '{safe_shape_name(shape)}'; description quality may be weak.",
                    "visual",
                    True,
                )
            )

        if description_source == "missing" and WARN_ON_MISSING_ALT_TEXT:
            visual_warnings.append(
                make_warning(
                    "DESCRIPTION_METADATA_MISSING",
                    f"Detected visual '{safe_shape_name(shape)}' has no description metadata.",
                    "visual",
                    True,
                )
            )

        visual_entry = {
            "observed_source": {
                "package_part": package_part,
                "relationships_part": relationships_part,
                "shape_tree_path": f"/p:sld/p:cSld/p:spTree/{safe_shape_kind(shape)}[{index}]",
                "shape_kind": safe_shape_kind(shape),
                "shape_id": safe_shape_id(shape),
                "shape_name": safe_shape_name(shape),
                "geometry": geometry_from_shape(shape),
                "related_part": None,
                "non_visual_properties": {
                    "descr": description if description_source == "descr" else "",
                    "title": description if description_source == "title" else "",
                },
                "from_group": ctx.inside_group,
            },
            "normalized_workflow": {
                "label": f"Visual {visual_index}",
                "visual_type": visual_type,
                "description": description,
                "description_source": description_source,
                "manual_review_required": True,
            },
            "projected_target": {
                "target_package_part": "word/document.xml",
                "target_object_kind": "visual_description_paragraph",
                "target_path_hint": "/w:document/w:body",
                "associated_drawing_reference": None,
                "projection_status": "planned",
            },
            "warning_entries": visual_warnings,
        }

        visual_entries.append(visual_entry)
        warnings.extend(visual_warnings)

    return visual_entries, dedupe_warning_entries(warnings)


def get_shape_alt_text(shape: Any) -> tuple[str, str]:
    try:
        for element in shape.element.iter():
            tag = str(getattr(element, "tag", ""))
            if tag.endswith("}cNvPr"):
                descr = normalize_paragraph_text(element.get("descr", ""))
                if descr:
                    return descr, "descr"
                title = normalize_paragraph_text(element.get("title", ""))
                if title:
                    return title, "title"
    except Exception:
        return "", "missing"
    return "", "missing"


def infer_visual_type(shape: Any, description_text: str) -> tuple[str, bool]:
    name_blob = " ".join(part for part in [safe_shape_name(shape), description_text] if part).lower()

    if getattr(shape, "has_chart", False):
        return "chart", False

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        if any(token in name_blob for token in ("screenshot", "screen shot", "ui", "dashboard", "interface")):
            return "screenshot", False
        if "table" in name_blob or "grid" in name_blob:
            return "table image", False
        if "diagram" in name_blob or "flow" in name_blob or "process" in name_blob:
            return "diagram", False
        if "icon" in name_blob or "icons" in name_blob:
            return "icon set", False
        return "photo", True

    return "other detected visual", True


def extract_text_block_entries(
    package_part: str,
    shape_contexts: list[ShapeContext],
    relationship_warning_present: bool,
    visual_shapes: set[int],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    text_block_entries: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []

    shape_index = 0
    for ctx in shape_contexts:
        if ctx.is_group_container:
            continue
        shape_index += 1
        shape = ctx.shape
        shape_tree_path = f"/p:sld/p:cSld/p:spTree/{safe_shape_kind(shape)}[{shape_index}]"

        try:
            if getattr(shape, "has_text_frame", False):
                paragraph_entries = extract_paragraph_entries(shape.text_frame, "p:txBody")
                if paragraph_entries:
                    near_visual = is_near_visual(shape, shape_contexts, visual_shapes)
                    entry = build_text_block_entry(
                        package_part,
                        shape_tree_path,
                        shape,
                        "text",
                        safe_shape_name(shape),
                        paragraph_entries,
                        ctx.inside_group,
                        near_visual,
                        relationship_warning_present,
                    )
                    text_block_entries.append(entry)

            if getattr(shape, "has_table", False):
                table_entries = extract_table_text_blocks(
                    package_part,
                    shape_tree_path,
                    shape,
                    ctx.inside_group,
                    relationship_warning_present,
                )
                text_block_entries.extend(table_entries)

        except Exception as exc:
            if WARN_ON_UNSUPPORTED_SHAPES:
                warnings.append(
                    make_warning(
                        "TEXT_EXTRACTION_PARTIAL_FAILURE",
                        f"Could not fully extract text from shape '{safe_shape_name(shape)}': {exc}",
                        "slide",
                        True,
                    )
                )

    text_block_entries = sort_text_block_entries(text_block_entries)
    for position, entry in enumerate(text_block_entries, start=1):
        entry["normalized_workflow"]["ordered_position"] = position

    return text_block_entries, dedupe_warning_entries(warnings)


def extract_slide_entries(
    presentation: Any,
    preview_lookup: dict[int, Path],
    preview_warnings_by_slide: dict[int, list[dict[str, Any]]],
) -> list[dict[str, Any]]:
    slide_entries: list[dict[str, Any]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        if SHOW_PROGRESS:
            LOGGER.info("Processing slide %s", slide_number)

        package_part = f"ppt/slides/slide{slide_number}.xml"
        relationships_part = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
        preview_path = preview_lookup.get(slide_number)
        preview_warnings = preview_warnings_by_slide.get(slide_number, [])
        shape_contexts = list(iter_shape_contexts(slide.shapes))

        visual_shape_ids = {
            id(ctx.shape)
            for ctx in shape_contexts
            if not ctx.is_group_container and (
                getattr(ctx.shape, "has_chart", False)
                or getattr(ctx.shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
            )
        }

        relationship_warning = detect_relationship_meaning_warning(shape_contexts)
        skipped_shape_warnings = detect_skipped_shape_warnings(shape_contexts)

        text_block_entries, text_block_warnings = extract_text_block_entries(
            package_part,
            shape_contexts,
            relationship_warning_present=relationship_warning is not None,
            visual_shapes=visual_shape_ids,
        )

        title, title_source, title_warnings = choose_title(slide, text_block_entries, slide_number)
        visual_entries, visual_warnings = extract_visual_entries(package_part, relationships_part, shape_contexts)
        slide_text_lines = flatten_text_block_entries_for_docx(text_block_entries)
        reading_order_low_confidence = is_reading_order_weak(text_block_entries)
        grouped_content_detected = any(ctx.is_group_container for ctx in shape_contexts)
        skipped_shape_content_detected = bool(skipped_shape_warnings)
        preview_match_ambiguous = bool(preview_warnings)

        slide_warning_entries: list[dict[str, Any]] = []
        slide_warning_entries.extend(title_warnings)
        slide_warning_entries.extend(text_block_warnings)
        slide_warning_entries.extend(visual_warnings)
        slide_warning_entries.extend(skipped_shape_warnings)
        slide_warning_entries.extend(preview_warnings)

        if grouped_content_detected and WARN_ON_UNSUPPORTED_SHAPES:
            slide_warning_entries.append(
                make_warning(
                    "GROUPED_CONTENT_DETECTED",
                    "Grouped content detected; extraction may be incomplete.",
                    "slide",
                    True,
                )
            )

        if relationship_warning is not None:
            slide_warning_entries.append(relationship_warning)

        if not slide_text_lines:
            slide_warning_entries.append(
                make_warning(
                    "NO_EXTRACTABLE_TEXT",
                    "No extractable visible slide text found on this slide.",
                    "slide",
                    True,
                )
            )

        if reading_order_low_confidence:
            slide_warning_entries.append(
                make_warning(
                    "READING_ORDER_LOW_CONFIDENCE",
                    "Reading-order heuristic may be weak for this slide.",
                    "slide",
                    True,
                )
            )

        if INCLUDE_PREVIEW_IMAGES and PREVIEW_IMAGE_DIR and preview_path is None:
            slide_warning_entries.append(
                make_warning(
                    "PREVIEW_NOT_FOUND",
                    "Preview image not found for this slide.",
                    "slide",
                    True,
                )
            )

        if visual_entries and len(text_block_entries) <= VISUAL_HEAVY_TEXT_BLOCK_THRESHOLD:
            slide_warning_entries.append(
                make_warning(
                    "IMAGE_HEAVY_RELATIVE_TO_TEXT",
                    "Slide appears image-heavy relative to extractable text; rendered meaning may be underrepresented.",
                    "slide",
                    True,
                )
            )

        slide_entries.append(
            {
                "slide_number": slide_number,
                "observed_source": {
                    "package_part": package_part,
                    "relationships_part": relationships_part,
                    "presentationml_root": "p:sld",
                    "common_slide_path": "/p:sld/p:cSld",
                    "shape_tree_path": "/p:sld/p:cSld/p:spTree",
                    "notes_part_present": False,
                    "preview_image_candidates": [relative_preview_candidate(preview_path)] if preview_path else [],
                },
                "normalized_workflow": {
                    "chosen_title": title,
                    "title_source": title_source,
                    "grouped_content_detected": grouped_content_detected,
                    "reading_order_low_confidence": reading_order_low_confidence,
                    "skipped_shape_content_detected": skipped_shape_content_detected,
                    "preview_match_ambiguous": preview_match_ambiguous,
                    "connector_relationship_warning": relationship_warning is not None,
                    "manual_review_required": True,
                },
                "projected_target": {
                    "target_package_part": "word/document.xml",
                    "projected_slide_heading_text": f"Slide {slide_number}. {title}",
                    "projected_slide_heading_style": "Heading1",
                    "projected_slide_separator_mode": SLIDE_SEPARATOR_MODE,
                    "projection_status": "planned",
                },
                "warning_entries": dedupe_warning_entries(slide_warning_entries),
                "text_block_entries": text_block_entries,
                "visual_entries": visual_entries,
            }
        )

    return slide_entries


def build_manifest(input_path: Path, output_paths: OutputPaths, slide_entries: list[dict[str, Any]]) -> dict[str, Any]:
    total_slides = len(slide_entries)
    slides_with_warnings = sum(1 for slide in slide_entries if slide["warning_entries"])
    visual_entries = sum(len(slide["visual_entries"]) for slide in slide_entries)
    visuals_missing_descriptions = sum(
        1
        for slide in slide_entries
        for visual in slide["visual_entries"]
        if visual["normalized_workflow"]["description_source"] == "missing"
    )

    return {
        "manifest_version": "1.0",
        "manifest_kind": "pptx_to_docx_remediation_manifest",
        "generated_at_utc": TIMESTAMP_UTC,
        "generator": {
            "script_name": "script_1_extract",
            "script_version": "1.0.0",
        },
        "source_package": {
            "format_family": "OOXML",
            "package_model": "OPC",
            "document_vocabulary": "PresentationML",
            "input_file_name": input_path.name,
            "input_file_path": str(input_path),
        },
        "target_package": {
            "format_family": "OOXML",
            "package_model": "OPC",
            "document_vocabulary": "WordprocessingML",
            "draft_output_file_name": output_paths.docx.name,
            "draft_output_file_path": str(output_paths.docx),
        },
        "run_config": {
            "include_preview_images": INCLUDE_PREVIEW_IMAGES,
            "preview_image_dir": PREVIEW_IMAGE_DIR or None,
            "title_fallback_mode": TITLE_FALLBACK_MODE,
            "text_order_mode": TEXT_ORDER_MODE,
            "slide_separator_mode": SLIDE_SEPARATOR_MODE,
        },
        "manifest_summary": {
            "total_slides": total_slides,
            "slides_with_warnings": slides_with_warnings,
            "detected_visual_entries": visual_entries,
            "detected_visual_entries_missing_descriptions": visuals_missing_descriptions,
        },
        "manifest_warning_entries": [],
        "slide_entries": slide_entries,
    }


def slide_image_alt_metadata(slide: dict[str, Any]) -> tuple[str, str]:
    slide_number = slide["slide_number"]
    chosen_title = slide["normalized_workflow"]["chosen_title"]
    title = f"Slide {slide_number} image: {chosen_title}"
    description = (
        f"Image of slide {slide_number}, titled '{chosen_title}'. "
        "The slide's extracted text and visual descriptions are provided as real Word content following this image."
    )
    return title, description


def apply_picture_alt_text(inline_shape: Any, title: str, description: str) -> None:
    doc_pr = inline_shape._inline.docPr
    doc_pr.set("title", title)
    doc_pr.set("descr", description)


def enforce_wcag_aaa_text_colors(document: Document) -> None:
    for style in document.styles:
        try:
            style.font.color.rgb = WCAG_AAA_TEXT_COLOR
        except Exception:
            pass

        try:
            color_elements = style.element.xpath(".//w:rPr/w:color")
        except Exception:
            color_elements = []
        for color_element in color_elements:
            color_element.set(qn("w:val"), WCAG_AAA_TEXT_COLOR_HEX)
            for attr_name in ("w:themeColor", "w:themeTint", "w:themeShade"):
                color_element.attrib.pop(qn(attr_name), None)

    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = WCAG_AAA_TEXT_COLOR


def create_docx_draft(manifest: dict[str, Any]) -> Document:
    title = DOC_TITLE.strip() if DOC_TITLE.strip() else f"{Path(manifest['source_package']['input_file_name']).stem} companion document"
    document = Document()
    document.core_properties.title = title
    document.add_heading(title, level=0)

    slide_entries = manifest["slide_entries"]
    for index, slide in enumerate(slide_entries, start=1):
        document.add_heading(slide["projected_target"]["projected_slide_heading_text"], level=1)

        preview_candidates = slide["observed_source"].get("preview_image_candidates", [])
        if INCLUDE_PREVIEW_IMAGES and preview_candidates and PREVIEW_IMAGE_DIR:
            preview_path = Path(PREVIEW_IMAGE_DIR) / preview_candidates[0]
            if preview_path.exists():
                try:
                    inline_shape = document.add_picture(str(preview_path), width=Inches(6.5))
                    alt_title, alt_description = slide_image_alt_metadata(slide)
                    apply_picture_alt_text(inline_shape, alt_title, alt_description)
                except Exception:
                    slide["warning_entries"].append(
                        make_warning(
                            "SLIDE_IMAGE_INSERTION_FAILED",
                            f"Slide image insertion failed for '{preview_path.name}'.",
                            "slide",
                            True,
                        )
                    )

        document.add_heading(SECTION_HEADING_SLIDE_TEXT, level=2)
        text_lines = flatten_text_block_entries_for_docx(slide["text_block_entries"])
        if text_lines:
            for line in text_lines:
                document.add_paragraph(line)
        elif INCLUDE_EMPTY_SLIDE_TEXT_SECTION:
            document.add_paragraph(PLACEHOLDER_MISSING_TEXT)

        if slide["visual_entries"] or INCLUDE_EMPTY_VISUAL_DESCRIPTIONS_SECTION:
            document.add_heading(SECTION_HEADING_VISUALS, level=2)
            if slide["visual_entries"]:
                for visual in slide["visual_entries"]:
                    paragraph = document.add_paragraph()
                    label = visual["normalized_workflow"]["label"]
                    visual_type = visual["normalized_workflow"]["visual_type"].capitalize()
                    description = visual["normalized_workflow"]["description"]
                    description_source = visual["normalized_workflow"]["description_source"]

                    paragraph.add_run(f"{label}. {visual_type}: ").bold = True
                    if description:
                        text = description
                        if description_source == "title":
                            text += " [metadata source: title]"
                    else:
                        text = PLACEHOLDER_MISSING_DESCRIPTION
                    paragraph.add_run(text)
            else:
                document.add_paragraph("No detected charts or pictures identified.")

        if index < len(slide_entries):
            if SLIDE_SEPARATOR_MODE == "PAGE_BREAK":
                p = document.add_paragraph()
                p.add_run().add_break(WD_BREAK.PAGE)
            else:
                document.add_paragraph("")

    enforce_wcag_aaa_text_colors(document)
    return document


def build_review_notes(manifest: dict[str, Any]) -> str:
    slides = manifest["slide_entries"]
    warning_slides = [slide for slide in slides if slide["warning_entries"]]

    lines = [
        "# Review notes",
        "",
        "## Summary",
        "",
        "- This tool extracts text from standard text containers and table cells.",
        "- It creates separate visual entries for detected charts and pictures, using existing description metadata when available and a placeholder when it is missing.",
        "- It does not preserve full layout meaning, diagram relationships, Word list structure, or accessibility structure.",
        f"- Total slides: {len(slides)}",
        f"- Slides with warnings: {len(warning_slides)}",
        f"- Detected visual entries missing descriptions: {manifest['manifest_summary']['detected_visual_entries_missing_descriptions']}",
        "",
    ]

    for slide in warning_slides:
        lines.append(f"## Slide {slide['slide_number']}. {slide['normalized_workflow']['chosen_title']}")
        lines.append("")
        for warning in slide["warning_entries"]:
            lines.append(f"- {warning['warning_code']} — {warning['warning_message']}")
        lines.append("")

    return "\n".join(lines).rstrip() + "\n"


def build_extract_report(manifest: dict[str, Any]) -> str:
    lines = [
        "# Extract Report",
        "",
        "## Manifest Summary",
        "",
        f"- Manifest version: {manifest['manifest_version']}",
        f"- Source vocabulary: {manifest['source_package']['document_vocabulary']}",
        f"- Target vocabulary: {manifest['target_package']['document_vocabulary']}",
        f"- Total slides: {manifest['manifest_summary']['total_slides']}",
        f"- Slides with warnings: {manifest['manifest_summary']['slides_with_warnings']}",
        f"- Detected visual entries: {manifest['manifest_summary']['detected_visual_entries']}",
        f"- Detected visual entries missing descriptions: {manifest['manifest_summary']['detected_visual_entries_missing_descriptions']}",
        "",
    ]

    for slide in manifest["slide_entries"]:
        lines.append(f"## Slide {slide['slide_number']}. {slide['normalized_workflow']['chosen_title']}")
        lines.append("")
        lines.append(f"**Title source:** {slide['normalized_workflow']['title_source']}")
        lines.append(f"**Grouped content detected:** {'Yes' if slide['normalized_workflow']['grouped_content_detected'] else 'No'}")
        lines.append(f"**Reading order low confidence:** {'Yes' if slide['normalized_workflow']['reading_order_low_confidence'] else 'No'}")
        lines.append(f"**Skipped shape content detected:** {'Yes' if slide['normalized_workflow']['skipped_shape_content_detected'] else 'No'}")
        lines.append(f"**Preview match ambiguous:** {'Yes' if slide['normalized_workflow']['preview_match_ambiguous'] else 'No'}")
        lines.append(f"**Connector relationship warning:** {'Yes' if slide['normalized_workflow']['connector_relationship_warning'] else 'No'}")
        lines.append(f"**Manual review required:** {'Yes' if slide['normalized_workflow']['manual_review_required'] else 'No'}")
        lines.append("")

        lines.append("### Slide-level warnings")
        lines.append("")
        if slide["warning_entries"]:
            for warning in slide["warning_entries"]:
                lines.append(f"- {warning['warning_code']} — {warning['warning_message']}")
        else:
            lines.append("- None")
        lines.append("")

        lines.append("### Extracted text blocks")
        lines.append("")
        if slide["text_block_entries"]:
            for index, block in enumerate(slide["text_block_entries"], start=1):
                observed = block["observed_source"]
                normalized = block["normalized_workflow"]
                lines.append(f"#### Text block {index}")
                lines.append("")
                lines.append(f"- Source kind: {normalized['source_kind']}")
                lines.append(f"- Source name: {normalized['source_name']}")
                lines.append(f"- Source text container type: {observed['source_text_container_type']}")
                lines.append(f"- Placeholder semantics: {'Yes' if observed['has_placeholder_semantics'] else 'No'}")
                lines.append(f"- Placeholder type: {observed['placeholder_type']}")
                lines.append(f"- Interpreted text role: {normalized['interpreted_text_role']}")
                lines.append(f"- Role confidence: {normalized['role_confidence']}")
                lines.append(f"- Auto-merge allowed: {'Yes' if normalized['auto_merge_allowed'] else 'No'}")
                lines.append("")
                if block["warning_entries"]:
                    lines.append("Warnings:")
                    for warning in block["warning_entries"]:
                        lines.append(f"- {warning['warning_code']} — {warning['warning_message']}")
                    lines.append("")
                lines.append("Paragraphs:")
                for paragraph in block["paragraph_entries"]:
                    pnorm = paragraph["normalized_workflow"]
                    lines.append(f"- [{pnorm['list_kind']}][level {pnorm['level']}] {pnorm['text']}")
                    for pwarn in paragraph["warning_entries"]:
                        lines.append(f"  - Warning: {pwarn['warning_code']} — {pwarn['warning_message']}")
                lines.append("")
        else:
            lines.append("- None")
            lines.append("")

        lines.append("### Visual entries")
        lines.append("")
        if slide["visual_entries"]:
            for visual in slide["visual_entries"]:
                vnorm = visual["normalized_workflow"]
                lines.append(f"#### {vnorm['label']}")
                lines.append("")
                lines.append(f"- Type: {vnorm['visual_type']}")
                lines.append(f"- Description source: {vnorm['description_source']}")
                lines.append(f"- Description: {vnorm['description'] or '[missing]'}")
                lines.append(f"- Manual review required: {'Yes' if vnorm['manual_review_required'] else 'No'}")
                if visual["warning_entries"]:
                    lines.append("")
                    lines.append("Warnings:")
                    for warning in visual["warning_entries"]:
                        lines.append(f"- {warning['warning_code']} — {warning['warning_message']}")
                lines.append("")
        else:
            lines.append("- None")
            lines.append("")

    return "\n".join(lines).rstrip() + "\n"


def atomic_write_text(path: Path, content: str) -> None:
    temp_path: Optional[Path] = None
    try:
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=path.suffix or ".tmp",
            dir=str(path.parent),
            prefix=f"{path.stem}_{datetime.now().strftime('%Y%m%d%H%M%S')}_",
        ) as temp_file:
            temp_path = Path(temp_file.name)

        temp_path.write_text(content, encoding="utf-8")
        temp_path.replace(path)
    except Exception as exc:
        if temp_path is not None and temp_path.exists():
            temp_path.unlink(missing_ok=True)
        raise RuntimeError(f"Failed to write file '{path}': {exc}") from exc


def atomic_write_docx(document: Document, path: Path) -> None:
    temp_path: Optional[Path] = None
    try:
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".docx",
            dir=str(path.parent),
            prefix=f"{path.stem}_{datetime.now().strftime('%Y%m%d%H%M%S')}_",
        ) as temp_file:
            temp_path = Path(temp_file.name)

        document.save(str(temp_path))
        temp_path.replace(path)
    except Exception as exc:
        if temp_path is not None and temp_path.exists():
            temp_path.unlink(missing_ok=True)
        raise RuntimeError(f"Failed to write DOCX '{path}': {exc}") from exc


def write_outputs(output_paths: OutputPaths, manifest: dict[str, Any]) -> tuple[bool, bool, bool, bool, bool]:
    if DRY_RUN:
        LOGGER.info("DRY_RUN: would write DOCX to %s", output_paths.docx)
        LOGGER.info("DRY_RUN: would write JSON manifest to %s", output_paths.manifest_json)
        LOGGER.info("DRY_RUN: would write YAML manifest to %s", output_paths.manifest_yaml)
        LOGGER.info("DRY_RUN: would write extract report to %s", output_paths.extract_report_md)
        LOGGER.info("DRY_RUN: would write review notes to %s", output_paths.review_notes_md)
        return False, False, False, False, False

    document = create_docx_draft(manifest)
    review_notes = build_review_notes(manifest)
    extract_report = build_extract_report(manifest)
    json_text = json.dumps(manifest, indent=2, ensure_ascii=False)
    yaml_text = yaml.safe_dump(manifest, sort_keys=False, allow_unicode=True)

    docx_written = json_written = yaml_written = report_written = notes_written = False

    atomic_write_docx(document, output_paths.docx)
    docx_written = True

    atomic_write_text(output_paths.manifest_json, json_text + "\n")
    json_written = True

    atomic_write_text(output_paths.manifest_yaml, yaml_text)
    yaml_written = True

    atomic_write_text(output_paths.extract_report_md, extract_report)
    report_written = True

    if WRITE_REVIEW_NOTES:
        atomic_write_text(output_paths.review_notes_md, review_notes)
        notes_written = True

    return docx_written, json_written, yaml_written, report_written, notes_written


def build_run_summary(
    status: str,
    message: str,
    input_path: Path,
    output_paths: OutputPaths,
    manifest: dict[str, Any],
    outputs_written: tuple[bool, bool, bool, bool, bool],
) -> dict[str, Any]:
    docx_written, json_written, yaml_written, report_written, notes_written = outputs_written
    return {
        "status": status,
        "message": message,
        "input_pptx": str(input_path),
        "output_docx": str(output_paths.docx),
        "output_json_manifest": str(output_paths.manifest_json),
        "output_yaml_manifest": str(output_paths.manifest_yaml),
        "output_extract_report": str(output_paths.extract_report_md),
        "output_review_notes": str(output_paths.review_notes_md),
        "docx_written": docx_written,
        "json_manifest_written": json_written,
        "yaml_manifest_written": yaml_written,
        "extract_report_written": report_written,
        "review_notes_written": notes_written,
        "total_slides": manifest["manifest_summary"]["total_slides"],
        "slides_with_warnings": manifest["manifest_summary"]["slides_with_warnings"],
        "detected_visual_entries": manifest["manifest_summary"]["detected_visual_entries"],
        "detected_visual_entries_missing_descriptions": manifest["manifest_summary"]["detected_visual_entries_missing_descriptions"],
    }


def print_run_summary(summary: dict[str, Any]) -> None:
    lines = [
        f"Run status: {summary['status']}",
        f"Message: {summary['message']}",
        f"Dry run: {DRY_RUN}",
        f"Input PPTX: {summary['input_pptx']}",
        f"Output DOCX: {summary['output_docx']}",
        f"Output JSON manifest: {summary['output_json_manifest']}",
        f"Output YAML manifest: {summary['output_yaml_manifest']}",
        f"Output extract report: {summary['output_extract_report']}",
        f"Output review notes: {summary['output_review_notes']}",
        f"DOCX written: {summary['docx_written'] if not DRY_RUN else 'would write'}",
        f"JSON manifest written: {summary['json_manifest_written'] if not DRY_RUN else 'would write'}",
        f"YAML manifest written: {summary['yaml_manifest_written'] if not DRY_RUN else 'would write'}",
        f"Extract report written: {summary['extract_report_written'] if not DRY_RUN else 'would write'}",
        f"Review notes written: {summary['review_notes_written'] if not DRY_RUN else 'would write'}",
        f"Slides processed: {summary['total_slides']}",
        f"Slides with warnings: {summary['slides_with_warnings']}",
        f"Detected visual entries: {summary['detected_visual_entries']}",
        f"Detected visual entries missing descriptions: {summary['detected_visual_entries_missing_descriptions']}",
    ]
    print("\n".join(lines))


def main() -> int:
    setup_logging()

    input_path = Path(INPUT_PATH).expanduser() if INPUT_PATH else Path("./input.pptx")
    output_paths = derive_output_paths(input_path)
    manifest: dict[str, Any] = {
        "manifest_version": "1.0",
        "manifest_kind": "pptx_to_docx_remediation_manifest",
        "generated_at_utc": TIMESTAMP_UTC,
        "generator": {"script_name": "script_1_extract", "script_version": "1.0.0"},
        "source_package": {},
        "target_package": {},
        "run_config": {},
        "manifest_summary": {
            "total_slides": 0,
            "slides_with_warnings": 0,
            "detected_visual_entries": 0,
            "detected_visual_entries_missing_descriptions": 0,
        },
        "manifest_warning_entries": [],
        "slide_entries": [],
    }

    try:
        input_path, output_paths, preview_dir = validate_config()
        ensure_output_directories(output_paths)

        preview_dir = render_slide_images_if_needed(input_path, output_paths, preview_dir)
        preview_lookup, preview_warnings_by_slide = discover_preview_images(preview_dir)
        presentation = load_presentation(input_path)
        slide_entries = extract_slide_entries(presentation, preview_lookup, preview_warnings_by_slide)
        manifest = build_manifest(input_path, output_paths, slide_entries)

        if DRY_RUN:
            outputs_written = (False, False, False, False, False)
            status = STATUS_DRY_RUN_COMPLETE
            message = "Validation, extraction, manifest construction, and reporting completed. No files were written."
        else:
            try:
                outputs_written = write_outputs(output_paths, manifest)
                status = (
                    STATUS_SUCCESS_WITH_WARNINGS
                    if manifest["manifest_summary"]["slides_with_warnings"] > 0
                    else STATUS_SUCCESS
                )
                message = (
                    "Extraction completed. Manual review is still required."
                    if status == STATUS_SUCCESS_WITH_WARNINGS
                    else "Extraction completed without extraction warnings."
                )
            except Exception as output_exc:
                outputs_written = (False, False, False, False, False)
                status = STATUS_FAILURE
                message = str(output_exc)

                summary = build_run_summary(status, message, input_path, output_paths, manifest, outputs_written)
                print_run_summary(summary)
                LOGGER.error(message)
                return EXIT_RUNTIME_ERROR

        summary = build_run_summary(status, message, input_path, output_paths, manifest, outputs_written)
        print_run_summary(summary)

        if status == STATUS_SUCCESS_WITH_WARNINGS:
            LOGGER.warning("Build completed with extraction warnings. Manual review remains required.")
        elif status == STATUS_SUCCESS:
            LOGGER.info("Build completed cleanly.")
        elif status == STATUS_DRY_RUN_COMPLETE:
            LOGGER.info("Dry run completed.")

        return EXIT_SUCCESS

    except ConfigError as exc:
        summary = build_run_summary(
            STATUS_FAILURE,
            str(exc),
            input_path,
            output_paths,
            manifest,
            (False, False, False, False, False),
        )
        print_run_summary(summary)
        return fail(str(exc), EXIT_CONFIG_ERROR)

    except Exception as exc:
        summary = build_run_summary(
            STATUS_FAILURE,
            str(exc),
            input_path,
            output_paths,
            manifest,
            (False, False, False, False, False),
        )
        print_run_summary(summary)
        return fail(str(exc), EXIT_RUNTIME_ERROR)


if __name__ == "__main__":
    sys.exit(main())
