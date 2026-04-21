from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from .errors import ExtractionError
from .models import Manifest, ParagraphEntry, SlideEntry, TextBlockEntry, VisualEntry, WarningEntry
from .rendering import RenderResult
from .text_utils import clean_text


@dataclass(frozen=True)
class ShapeContext:
    shape: Any
    inside_group: bool = False


def extract_manifest(pptx_path: Path, output_docx_path: Path, render_result: RenderResult) -> Manifest:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise ExtractionError("python-pptx is required for extraction.") from exc

    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:
        raise ExtractionError(f"Failed to open PPTX '{pptx_path}': {exc}") from exc

    slides: list[SlideEntry] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        contexts = list(iter_shape_contexts(slide.shapes))
        text_blocks = extract_text_blocks(contexts)
        visuals = extract_visuals(contexts)
        title, title_source, title_warnings = choose_title(text_blocks, slide_number)
        slide_warnings = list(title_warnings)
        slide_warnings.extend(detect_slide_warnings(contexts, text_blocks, visuals))
        preview_image = None
        if slide_number in render_result.images_by_slide:
            preview_image = render_result.images_by_slide[slide_number].name
        elif render_result.image_dir is not None:
            slide_warnings.append(
                WarningEntry(
                    code="SLIDE_IMAGE_MISSING",
                    message="Slide image rendering was enabled, but no image was produced for this slide.",
                    scope="slide",
                )
            )

        slides.append(
            SlideEntry(
                slide_number=slide_number,
                title=title,
                title_source=title_source,
                preview_image=preview_image,
                text_blocks=tuple(text_blocks),
                visuals=tuple(visuals),
                warnings=tuple(dedupe_warnings(slide_warnings)),
            )
        )

    return Manifest(
        source_path=pptx_path,
        output_docx_path=output_docx_path,
        preview_image_dir=render_result.image_dir,
        slides=tuple(slides),
        warnings=render_result.warnings,
    )


def iter_shape_contexts(shapes: Iterable[Any], inside_group: bool = False) -> Iterable[ShapeContext]:
    for shape in shapes:
        yield ShapeContext(shape=shape, inside_group=inside_group)
        if safe_shape_type_name(shape).upper() == "GROUP":
            try:
                yield from iter_shape_contexts(shape.shapes, inside_group=True)
            except Exception:
                continue


def extract_text_blocks(contexts: list[ShapeContext]) -> list[TextBlockEntry]:
    blocks: list[TextBlockEntry] = []
    for context in contexts:
        shape = context.shape
        if has_table(shape):
            blocks.extend(extract_table_blocks(shape, context.inside_group, len(blocks) + 1))
            continue

        if not getattr(shape, "has_text_frame", False):
            continue
        paragraphs = extract_paragraphs(shape.text_frame.paragraphs)
        if not paragraphs:
            continue

        placeholder_type = placeholder_kind(shape)
        source_type = "placeholder_text" if placeholder_type else "standalone_text_box"
        role = interpreted_role(placeholder_type)
        confidence = "high" if placeholder_type else "low"
        warnings: list[WarningEntry] = []
        if context.inside_group:
            warnings.append(WarningEntry("GROUPED_TEXT_CONTENT", "Text was inside grouped content.", "text_block"))
        if source_type == "standalone_text_box":
            warnings.append(
                WarningEntry(
                    "STANDALONE_TEXT_BOX_ROLE_UNCERTAIN",
                    "Standalone text box role is uncertain and should be reviewed before merging into surrounding content.",
                    "text_block",
                )
            )

        blocks.append(
            TextBlockEntry(
                source_name=safe_shape_name(shape),
                source_text_container_type=source_type,
                placeholder_type=placeholder_type,
                interpreted_text_role=role,
                role_confidence=confidence,
                auto_merge_allowed=False,
                order=len(blocks) + 1,
                paragraphs=tuple(paragraphs),
                warnings=tuple(warnings),
            )
        )

    return sorted(blocks, key=lambda block: block.order)


def extract_table_blocks(shape: Any, inside_group: bool, start_order: int) -> list[TextBlockEntry]:
    blocks: list[TextBlockEntry] = []
    try:
        rows = shape.table.rows
    except Exception:
        return blocks

    order = start_order
    for row_index, row in enumerate(rows, start=1):
        for col_index, cell in enumerate(row.cells, start=1):
            text = clean_text(getattr(cell, "text", ""))
            if not text:
                continue
            warnings = [WarningEntry("TABLE_CELL_TEXT_EXTRACTED", "PowerPoint table-cell text was extracted as a text block.", "text_block")]
            if inside_group:
                warnings.append(WarningEntry("GROUPED_TEXT_CONTENT", "Table text was inside grouped content.", "text_block"))
            blocks.append(
                TextBlockEntry(
                    source_name=f"{safe_shape_name(shape)} cell {row_index},{col_index}",
                    source_text_container_type="table_cell",
                    placeholder_type=None,
                    interpreted_text_role="body_content",
                    role_confidence="medium",
                    auto_merge_allowed=False,
                    order=order,
                    paragraphs=(ParagraphEntry(text=text),),
                    warnings=tuple(warnings),
                )
            )
            order += 1
    return blocks


def extract_paragraphs(paragraphs: Iterable[Any]) -> list[ParagraphEntry]:
    entries: list[ParagraphEntry] = []
    for paragraph in paragraphs:
        text = clean_text(getattr(paragraph, "text", ""))
        if not text:
            continue
        level = int(getattr(paragraph, "level", 0) or 0)
        list_kind, marker = infer_list_kind(paragraph)
        warnings: list[WarningEntry] = []
        if list_kind == "number":
            warnings.append(
                WarningEntry(
                    "NUMBERED_SEQUENCE_REVIEW",
                    "Numbering intent was observed, but ordered sequence should be reviewed in the DOCX.",
                    "paragraph",
                )
            )
        entries.append(ParagraphEntry(text=text, level=level, list_kind=list_kind, list_marker_hint=marker, warnings=tuple(warnings)))
    return entries


def infer_list_kind(paragraph: Any) -> tuple[str, str]:
    try:
        xml = paragraph._p.xml
    except Exception:
        xml = ""
    if "<a:buAutoNum" in xml:
        return "number", ""
    if "<a:buChar" in xml or "<a:buBlip" in xml:
        return "bullet", "•"
    return "plain", ""


def extract_visuals(contexts: list[ShapeContext]) -> list[VisualEntry]:
    visuals: list[VisualEntry] = []
    for context in contexts:
        shape = context.shape
        if not is_visual(shape):
            continue
        description, source = shape_alt_text(shape)
        warnings: list[WarningEntry] = []
        if source == "missing":
            warnings.append(WarningEntry("DESCRIPTION_METADATA_MISSING", f"Visual '{safe_shape_name(shape)}' has no description metadata.", "visual"))
        elif source == "title":
            warnings.append(WarningEntry("DESCRIPTION_SOURCE_TITLE_ONLY", f"Visual '{safe_shape_name(shape)}' only has title metadata.", "visual"))
        if context.inside_group:
            warnings.append(WarningEntry("GROUPED_VISUAL_CONTENT", "Visual was inside grouped content.", "visual"))
        visuals.append(
            VisualEntry(
                label=f"Visual {len(visuals) + 1}",
                visual_type=visual_type(shape, description),
                description=description,
                description_source=source,
                warnings=tuple(warnings),
            )
        )
    return visuals


def choose_title(text_blocks: list[TextBlockEntry], slide_number: int) -> tuple[str, str, list[WarningEntry]]:
    for block in text_blocks:
        if block.placeholder_type == "title" and block.paragraphs:
            return block.paragraphs[0].text, "TITLE_PLACEHOLDER", []
    for block in text_blocks:
        if block.paragraphs:
            return block.paragraphs[0].text, "FIRST_TEXT_LINE", [
                WarningEntry("NO_TITLE_PLACEHOLDER_FOUND", "No title placeholder was found; first text line was used as title.", "slide")
            ]
    return f"Untitled slide {slide_number}", "UNTITLED_FALLBACK", [
        WarningEntry("NO_EXTRACTABLE_TEXT", "No extractable visible slide text was found.", "slide")
    ]


def detect_slide_warnings(contexts: list[ShapeContext], text_blocks: list[TextBlockEntry], visuals: list[VisualEntry]) -> list[WarningEntry]:
    warnings: list[WarningEntry] = []
    if any(context.inside_group for context in contexts):
        warnings.append(WarningEntry("GROUPED_CONTENT_DETECTED", "Grouped content may hide structure or reading order.", "slide"))
    if any(is_connector_like(context.shape) for context in contexts):
        warnings.append(WarningEntry("CONNECTOR_RELATIONSHIP_WARNING", "Connectors or lines may convey meaning not preserved in text.", "slide"))
    if visuals and len(text_blocks) <= 1:
        warnings.append(WarningEntry("IMAGE_HEAVY_SLIDE_REVIEW", "Slide appears image-heavy; confirm important meaning is captured in text.", "slide"))
    return warnings


def placeholder_kind(shape: Any) -> str | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        name = shape.placeholder_format.type.name.lower()
    except Exception:
        return "other"
    if "title" in name:
        return "title"
    if "subtitle" in name:
        return "subtitle"
    if "body" in name or "object" in name:
        return "body"
    return "other"


def interpreted_role(placeholder_type: str | None) -> str:
    if placeholder_type == "title":
        return "title_candidate"
    if placeholder_type in {"body", "subtitle"}:
        return "body_content"
    return "unknown"


def shape_alt_text(shape: Any) -> tuple[str, str]:
    try:
        for element in shape.element.iter():
            if str(getattr(element, "tag", "")).endswith("}cNvPr"):
                descr = clean_text(element.get("descr", ""))
                if descr:
                    return descr, "descr"
                title = clean_text(element.get("title", ""))
                if title:
                    return title, "title"
    except Exception:
        pass
    return "", "missing"


def is_visual(shape: Any) -> bool:
    return bool(getattr(shape, "has_chart", False) or safe_shape_type_name(shape).upper() == "PICTURE")


def visual_type(shape: Any, description: str) -> str:
    if getattr(shape, "has_chart", False):
        return "chart"
    blob = f"{safe_shape_name(shape)} {description}".lower()
    if any(token in blob for token in ("screenshot", "screen shot", "interface", "dashboard")):
        return "screenshot"
    if any(token in blob for token in ("diagram", "flow", "process", "model")):
        return "diagram"
    return "picture"


def is_connector_like(shape: Any) -> bool:
    name = safe_shape_type_name(shape).upper()
    return any(token in name for token in ("CONNECTOR", "LINE", "FREEFORM"))


def has_table(shape: Any) -> bool:
    return bool(getattr(shape, "has_table", False))


def safe_shape_name(shape: Any) -> str:
    return clean_text(getattr(shape, "name", "") or getattr(shape, "shape_id", "") or "Unnamed shape")


def safe_shape_type_name(shape: Any) -> str:
    try:
        shape_type = getattr(shape, "shape_type", None)
        return getattr(shape_type, "name", str(shape_type or "UNKNOWN"))
    except Exception:
        return "UNKNOWN"


def dedupe_warnings(warnings: list[WarningEntry]) -> list[WarningEntry]:
    seen: set[tuple[str, str, str]] = set()
    output: list[WarningEntry] = []
    for warning in warnings:
        key = (warning.code, warning.message, warning.scope)
        if key in seen:
            continue
        seen.add(key)
        output.append(warning)
    return output
